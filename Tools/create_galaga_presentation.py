#!/usr/bin/env python3
"""
Create a PowerPoint presentation about the Galaga Wars project for Grade 10.
Uses website theme (dark bg, neon cyan/yellow), galagawars_logo, and sprite images.
Requires: pip install python-pptx
Usage: python create_galaga_presentation.py
Output: Galaga_Wars_Grade10_Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Website theme colors and font (from Website/index.html)
DARK_BG = (10, 14, 39)        # #0a0e27
DARK_BG_2 = (26, 31, 58)      # #1a1f3a
NEON_CYAN = (0, 255, 255)     # #00FFFF
NEON_YELLOW = (255, 255, 0)   # #FFFF00
NEON_MAGENTA = (255, 0, 255)  # #FF00FF
NEON_GREEN = (0, 255, 0)      # #00FF00
BODY_LIGHT = (200, 220, 255)  # light blue-gray for readability
WEBSITE_FONT = "VT323"        # same as website
BORDER_MARGIN_IN = 0.2         # inset for border (like website padding)
BORDER_WIDTH_PT = 3           # 3pt like website 3px

# Paths relative to project root (where script is run from)
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_PATH = os.path.join(PROJECT_ROOT, "sprites", "galagawars_logo", "2ed7f732-6374-4867-9c70-46be91f5c1ab.png")
SHUTTLE_DIR = os.path.join(PROJECT_ROOT, "sprites", "sImperialShuttle")
# 12 shuttle frames (different rotation angles) - smaller size on slide
SHUTTLE_FRAMES = [
    "93f461c0-d037-47ea-be01-a6709ab98196.png",
    "699ee271-fc5f-4328-afb6-a99320d645c6.png",
    "370f75fd-4364-4931-b19d-9d2a455a557e.png",
    "e7d18a59-aab7-49e9-bdff-eb6e62a65fb1.png",
    "0487de4f-ba82-4ef5-806a-94fcac59101e.png",
    "06fce25d-b629-41a1-bd28-16a701981a52.png",
    "086fbb8e-5da4-43b7-87d9-2a06070ee534.png",
    "16d3fba6-c564-4483-a1fe-cec7d3fbf43d.png",
    "17a42fcd-dcf3-4c40-b45b-5b9e55f8e977.png",
    "18c3e5e0-eeee-4de3-9e5f-c06207724a39.png",
    "29c70def-6900-4920-aae9-2313d32980d6.png",
    "2db613e8-ed51-4b17-acd8-9a289550c0cc.png",
]


def rgb(r, g, b):
    """Return RGBColor for (r,g,b) 0-255."""
    return RGBColor(r, g, b)


def apply_theme(slide):
    """Apply website-style dark background, neon text colors, and VT323 font."""
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(*DARK_BG)
    except Exception:
        pass
    # Title – neon cyan, VT323 (same as website)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        for p in tf.paragraphs:
            p.font.color.rgb = rgb(*NEON_CYAN)
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = WEBSITE_FONT
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = WEBSITE_FONT
    # Body placeholder – light text, VT323
    if len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        if ph.has_text_frame:
            for p in ph.text_frame.paragraphs:
                p.font.color.rgb = rgb(*BODY_LIGHT)
                p.font.size = Pt(18)
                p.font.name = WEBSITE_FONT
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = rgb(*BODY_LIGHT)
                    run.font.name = WEBSITE_FONT


def add_slide_border(slide, prs):
    """Add a border around the slide like the website (cyan + magenta top)."""
    try:
        margin = Inches(BORDER_MARGIN_IN)
        left = margin
        top = margin
        width = prs.slide_width - 2 * margin
        height = prs.slide_height - 2 * margin
        # Full rectangle border – cyan (like website bottom/sides)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.background()
        shape.line.color.rgb = rgb(*NEON_CYAN)
        shape.line.width = Pt(BORDER_WIDTH_PT)
        # Top accent – magenta bar (like website border-top)
        bar_height = Pt(BORDER_WIDTH_PT)
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, bar_height
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = rgb(*NEON_MAGENTA)
        top_bar.line.fill.background()
        # Send borders to back so content is on top (first shape = backmost)
        first_el = slide.shapes[0]._element
        first_el.addprevious(top_bar._element)
        first_el.addprevious(shape._element)
    except Exception:
        pass


def add_picture(slide, path, left_in, top_in, width_in=None):
    """Add picture to slide if file exists. Size by width if given."""
    if not path or not os.path.isfile(path):
        return None
    try:
        kw = {"width": Inches(width_in)} if width_in else {}
        return slide.shapes.add_picture(path, Inches(left_in), Inches(top_in), **kw)
    except Exception:
        return None


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide; returns slide for adding logo."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    apply_theme(slide)
    # Add Galaga Wars logo below title (centered)
    if os.path.isfile(LOGO_PATH):
        try:
            pic = slide.shapes.add_picture(LOGO_PATH, Inches(2.5), Inches(2.8), width=Inches(5))
            pic.left = int((prs.slide_width - pic.width) / 2)
        except Exception:
            pass
    add_slide_border(slide, prs)
    return slide


def add_content_slide(prs, title, bullet_points, notes="", image_path=None, image_left=None, image_top=None, image_width=None):
    """Add a content slide with bullet points. Optionally add one image."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    apply_theme(slide)
    # Re-apply body color after clearing
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.color.rgb = rgb(*BODY_LIGHT)
    if image_path and image_left is not None and image_top is not None:
        add_picture(slide, image_path, image_left, image_top, image_width)
    add_slide_border(slide, prs)
    return slide


def add_slide_with_sprite_grid(prs, title, bullet_points, sprite_paths, notes=""):
    """Add content slide with 12 sprite images in 2 rows of 6 (smaller size)."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    apply_theme(slide)
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.color.rgb = rgb(*BODY_LIGHT)
        p.font.name = WEBSITE_FONT
    # 12 shuttle sprites in 2 rows of 6, smaller size
    if sprite_paths:
        n = min(12, len(sprite_paths))
        cols = 6
        rows = (n + cols - 1) // cols
        w = 0.72   # inches per sprite (smaller)
        gap = 0.08
        total_w = cols * w + (cols - 1) * gap
        left_start = (10 - total_w) / 2
        top_row0 = 1.35
        row_height = w + gap
        for i, path in enumerate(sprite_paths[:n]):
            if os.path.isfile(path):
                try:
                    row, col = i // cols, i % cols
                    slide.shapes.add_picture(
                        path,
                        Inches(left_start + col * (w + gap)),
                        Inches(top_row0 + row * row_height),
                        width=Inches(w),
                    )
                except Exception:
                    pass
    add_slide_border(slide, prs)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title with logo ---
    add_title_slide(
        prs,
        "Galaga Wars: Game Design & Programming",
        "A Grade 10 Overview of a Real Game Project\nHow games are built with GameMaker Studio"
    )

    # --- Slide 2: What is GameMaker Studio? ---
    add_content_slide(
        prs,
        "What is GameMaker Studio?",
        [
            "GameMaker Studio is a 2D game engine used by indie and professional developers.",
            "You design your game using:",
            "  • Sprites (images), Objects (things in the game), Rooms (levels/screens)",
            "  • Scripts (code) written in GML (GameMaker Language)",
            "  • Events: Create, Step, Draw, Collision, etc.",
            "Games can be exported to: Windows, Mac, HTML5 (browser), consoles.",
            "Used for games like: Hyper Light Drifter, Undertale, and many mobile games.",
        ],
        notes="Brief intro to the tool - Grade 10 may not know it."
    )

    # --- Slide 3: Why GameMaker for This Project? ---
    add_content_slide(
        prs,
        "Why GameMaker for Galaga Wars?",
        [
            "Perfect for 2D arcade-style games (like the original Galaga from 1981).",
            "Visual editor + code: design levels and logic without only typing code.",
            "Built-in features: collisions, sprites, sounds, rooms.",
            "Export to web (HTML5) so anyone can play in a browser.",
            "Same concepts (sprites, objects, state machines) appear in bigger engines like Unity.",
        ]
    )

    # --- Slide 4: Project Structure ---
    add_content_slide(
        prs,
        "Galaga Wars – Project Structure",
        [
            "sprites/     – All images (ships, explosions, UI). E.g. sImperialShuttle, sTieFighter.",
            "objects/     – Things in the game: oPlayer, oTieFighter, oGameManager, oTitleScreen.",
            "scripts/     – Reusable code: WaveSpawner, ScoreManager, ObjectPool, GameConstants.",
            "rooms/       – Screens: SplashScreen, TitleScreen, GalagaWars (main game).",
            "datafiles/Patterns/ – JSON config: wave spawns, formation data, game_config.",
            "shaders/     – Visual effects (e.g. CRT scanline effect).",
            "sounds/      – Music and sound effects.",
        ],
        notes="Walk through folder structure so they see how a real project is organized."
    )

    # --- Slide 5: Key Elements of the Game ---
    add_content_slide(
        prs,
        "Key Elements of Galaga Wars",
        [
            "Game flow: Splash → Title → Play → Results → High Score entry.",
            "Enemies: TIE Fighter, TIE Interceptor, Imperial Shuttle (each with different behaviour).",
            "Formation: 5×8 grid of enemies; they dive, loop, and attack.",
            "Challenge stages: Bonus rounds every 4 levels.",
            "Tractor beam: Imperial Shuttle can capture the player ship.",
            "Score & lives: Extra lives at 20,000 and every 70,000 points.",
        ]
    )

    # --- Slide 6: 2D Sprites and "3D" Rotation – The Idea ---
    add_content_slide(
        prs,
        "2D Sprites That Look Like 3D Rotation",
        [
            "Games like this are 2D: everything is flat images on the screen.",
            "To make a ship look like it's rotating, we don't use real 3D – we use many 2D images.",
            "Each image is one \"angle\" of the ship (like a flipbook).",
            "Example: sImperialShuttle is one sprite made of 24 separate images.",
            "As the ship moves or turns, the game switches which image (frame) is shown.",
            "24 frames = 24 different angles → smooth rotation effect (15° per frame).",
        ],
        notes="This is the key concept: 2D art used to fake 3D rotation."
    )

    # --- Slide 7: sImperialShuttle – 24 Frames (with 4 sprite images) ---
    shuttle_paths = [os.path.join(SHUTTLE_DIR, f) for f in SHUTTLE_FRAMES]
    add_slide_with_sprite_grid(
        prs,
        "sImperialShuttle: 24 Frames for One Rotation",
        [
            "Sprite name: sImperialShuttle (Imperial Shuttle enemy).",
            "Total frames: 24 images in one sprite. Above: 12 sample frames (different angles).",
            "Each frame = one view of the shuttle (e.g. every 15° around a circle).",
            "In GameMaker: one \"sprite\" can have many \"sub-images\" (frames).",
            "Code uses image_index or image_speed to show the right frame.",
            "Same idea used for: sTieFighter, sTieIntercepter (multiple frames each).",
            "Result: on screen the ship appears to rotate smoothly, but it's just 2D art.",
        ],
        sprite_paths=shuttle_paths,
        notes="Reference: sprites/sImperialShuttle/ has 24 PNGs; sequence length 24, playback speed 8."
    )

    # --- Slide 8: How the Game Is Coded – Big Ideas ---
    add_content_slide(
        prs,
        "How the Game Is Coded – Big Ideas",
        [
            "State machines: Game has \"modes\" (title, playing, results). Enemies have states (formation, diving, looping).",
            "Data-driven design: Levels and behaviour come from JSON files, not only from code.",
            "Object pooling: Instead of creating/destroying bullets and explosions constantly, we reuse them (better performance).",
            "Managers: One object (e.g. oGameManager) coordinates others (spawning, score, challenges).",
            "Events: Create (when object appears), Step (every frame), Draw (every frame draw), Collision (when two things hit).",
        ]
    )

    # --- Slide 9: Code Structure Example ---
    add_content_slide(
        prs,
        "Code Structure – Who Does What",
        [
            "oGlobal (Create): Loads JSON, sets up global.Game, creates managers.",
            "oGameManager: Runs the main loop; switches game state (attract, play, results); calls WaveSpawner, ScoreManager.",
            "WaveSpawner: Reads wave_spawn.json; spawns the right enemies at the right positions.",
            "Enemy objects (e.g. oImperialShuttle): Have a state (formation / dive / loop); change sprite frame for rotation.",
            "oPlayer: Handles input, firing, collisions with enemies and bullets.",
        ],
        notes="No need to show actual code unless class is comfortable; focus on roles."
    )

    # --- Slide 10: What It Takes to Create a Polished Game ---
    add_content_slide(
        prs,
        "What It Takes to Create a Polished Game",
        [
            "Planning: Design documents, lists of features, and a clear \"loop\" (what the player does again and again).",
            "Art & audio: Many sprites (e.g. 24 per rotating ship), animations, sounds, music.",
            "Code quality: Organized scripts, clear names, state machines so behaviour is predictable.",
            "Data: Config files (JSON) so designers can change waves/difficulty without touching code.",
            "Testing: Play often; fix bugs; make sure it runs smoothly (e.g. 60 FPS).",
            "Polish: Effects (explosions, screen shake, CRT effect), feedback (sounds, score pop), and clear UI.",
        ]
    )

    # --- Slide 11: Polished Game – Summary ---
    add_content_slide(
        prs,
        "Polished Game – Summary",
        [
            "Consistent quality: Menus, gameplay, and \"game over\" all feel part of the same game.",
            "Performance: Object pooling and simple 2D graphics help keep the game smooth.",
            "Feedback: Player always knows what's happening (score, lives, sounds, visual effects).",
            "This project: 11,000+ lines of code, 40+ scripts, 35+ objects, 100+ assets – all for one arcade-style game.",
            "Takeaway: Even \"simple\" games need structure, data, and lots of small details to feel polished.",
        ]
    )

    # --- Slide 12: Questions & Try It ---
    add_title_slide(
        prs,
        "Questions?",
        "Try the game: run from GameMaker (F5) or play the web version.\nLook at sprites/sImperialShuttle to see the 24 rotation frames."
    )

    # Save
    out_path = os.path.join(PROJECT_ROOT, "Galaga_Wars_Grade10_Presentation.pptx")
    prs.save(out_path)
    print(f"Created: {out_path}")


if __name__ == "__main__":
    main()
